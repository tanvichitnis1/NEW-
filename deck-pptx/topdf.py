import base64, os, html as H
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PR = Presentation('EOV-Credentials-2026.pptx')
SW = PR.slide_width; SH = PR.slide_height
PXW, PXH = 1280, 720
sx = PXW / SW; sy = PXH / SH

def col(c, default=None):
    try:
        if c and c.type is not None and hasattr(c, 'rgb'):
            return '#%s' % str(c.rgb)
    except Exception:
        pass
    return default

def bg_of(sl):
    try:
        f = sl.background.fill
        if f.type is not None and f.fore_color and hasattr(f.fore_color, 'rgb'):
            return '#%s' % str(f.fore_color.rgb)
    except Exception:
        pass
    return '#FFFFFF'

ALIGN = {PP_ALIGN.CENTER:'center', PP_ALIGN.RIGHT:'right', PP_ALIGN.JUSTIFY:'justify'}
ANCH  = {MSO_ANCHOR.MIDDLE:'center', MSO_ANCHOR.BOTTOM:'flex-end'}

parts = []
for sl in PR.slides:
    body = []
    for sh in sl.shapes:
        if sh.left is None: continue
        x, y = sh.left*sx, sh.top*sy
        w, h = sh.width*sx, sh.height*sy
        base = f'position:absolute;left:{x:.2f}px;top:{y:.2f}px;width:{w:.2f}px;height:{h:.2f}px;'

        if sh.shape_type == 13 or sh.__class__.__name__ == 'Picture':          # picture
            blob = sh.image.blob
            b64 = base64.b64encode(blob).decode()
            body.append(f'<img style="{base}object-fit:contain" src="data:{sh.image.content_type};base64,{b64}">')
            continue

        # shape fill
        fill = None; alpha = 1.0; radius = 0
        try:
            if sh.fill.type is not None and sh.fill.type == 1:
                fill = col(sh.fill.fore_color)
                el = sh.fill._xPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                if el is not None: alpha = int(el.get('val'))/100000.0
        except Exception: pass
        if 'roundRect' in str(getattr(sh, 'shape_type', '')) or (sh.shape_type is not None and 'ROUND' in str(sh.shape_type)):
            radius = 10
        if fill:
            body.append(f'<div style="{base}background:{fill};opacity:{alpha:.3f};border-radius:{radius}px"></div>')

        if not sh.has_text_frame: continue
        tf = sh.text_frame
        if not tf.text.strip(): continue
        anchor = ANCH.get(tf.vertical_anchor, 'flex-start')
        inner = []
        for para in tf.paragraphs:
            if not ''.join(r.text for r in para.runs).strip():
                inner.append('<div style="height:.4em"></div>'); continue
            al = ALIGN.get(para.alignment, 'left')
            ls = ''
            try:
                if para.line_spacing and para.line_spacing > 3:
                    ls = f'line-height:{para.line_spacing.pt*(PXW/ (SW/914400) /72):.2f}px;'
            except Exception: pass
            runs = []
            for r in para.runs:
                f = r.font
                fs = (f.size.pt if f.size else 14) * (PXW/(SW/914400))/72
                c  = col(f.color, '#111111')
                cs = ''
                try:
                    el = f._rPr
                    if el is not None and el.get('spc'):
                        cs = f'letter-spacing:{int(el.get("spc"))/100*(PXW/(SW/914400))/72:.2f}px;'
                except Exception: pass
                runs.append(f'<span style="font-size:{fs:.2f}px;font-weight:{700 if f.bold else 400};color:{c};{cs}">{H.escape(r.text)}</span>')
            inner.append(f'<div style="text-align:{al};{ls}">{"".join(runs)}</div>')
        body.append(f'<div style="{base}display:flex;flex-direction:column;justify-content:{anchor};'
                    f'font-family:Arial,Helvetica,Arimo,sans-serif;overflow:hidden">{"".join(inner)}</div>')

    parts.append(f'<section style="background:{bg_of(sl)}">{"".join(body)}</section>')

doc = f'''<!doctype html><html><head><meta charset="utf-8">
<link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Arimo:wght@400;700&display=swap">
<style>
  @page {{ size: {PXW}px {PXH}px; margin: 0; }}
  html,body {{ margin:0; padding:0; }}
  section {{ position:relative; width:{PXW}px; height:{PXH}px; overflow:hidden;
            page-break-after:always; break-after:page; }}
  section:last-child {{ page-break-after:auto; }}
</style></head><body>{''.join(parts)}</body></html>'''
open('deck.html','w').write(doc)
print('wrote deck.html ·', len(PR.slides), 'slides')
